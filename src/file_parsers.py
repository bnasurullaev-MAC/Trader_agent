"""
Advanced file parsers for PPTX, DOCX, and PDF files.
"""

from __future__ import annotations
import logging
from pathlib import Path
from typing import Any, Dict, List, Optional, Tuple

logger = logging.getLogger(__name__)

def parse_pdf_file(pdf_path: Path) -> str:
    """
    Parse a PDF file and extract text content.
    
    Args:
        pdf_path: Path to the PDF file
        
    Returns:
        Extracted text content
    """
    try:
        import pdfplumber
        
        text_content = []
        
        with pdfplumber.open(pdf_path) as pdf:
            for page_num, page in enumerate(pdf.pages, 1):
                try:
                    # Extract text from the page
                    page_text = page.extract_text()
                    if page_text:
                        text_content.append(f"--- Page {page_num} ---\n{page_text}")
                    
                    # Extract tables if any
                    tables = page.extract_tables()
                    if tables:
                        for table_num, table in enumerate(tables, 1):
                            text_content.append(f"\n--- Table {table_num} on Page {page_num} ---")
                            for row in table:
                                if row:
                                    text_content.append(" | ".join(str(cell) if cell else "" for cell in row))
                
                except Exception as e:
                    logger.warning(f"Error processing page {page_num} of {pdf_path}: {e}")
                    continue
        
        return "\n\n".join(text_content)
        
    except ImportError:
        logger.error("pdfplumber not installed. Install with: pip install pdfplumber")
        return ""
    except Exception as e:
        logger.error(f"Error parsing PDF {pdf_path}: {e}")
        return ""

def parse_pptx_file(pptx_path: Path) -> str:
    """
    Parse a PPTX file and extract text content.
    
    Args:
        pptx_path: Path to the PPTX file
        
    Returns:
        Extracted text content
    """
    try:
        from pptx import Presentation
        
        text_content = []
        prs = Presentation(pptx_path)
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = []
            slide_text.append(f"--- Slide {slide_num} ---")
            
            # Extract text from all shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
                
                # Handle tables
                if shape.has_table:
                    table = shape.table
                    table_text = []
                    for row in table.rows:
                        row_text = []
                        for cell in row.cells:
                            if cell.text.strip():
                                row_text.append(cell.text.strip())
                        if row_text:
                            table_text.append(" | ".join(row_text))
                    if table_text:
                        slide_text.append("Table:")
                        slide_text.extend(table_text)
            
            if len(slide_text) > 1:  # More than just the slide header
                text_content.append("\n".join(slide_text))
        
        return "\n\n".join(text_content)
        
    except ImportError:
        logger.error("python-pptx not installed. Install with: pip install python-pptx")
        return ""
    except Exception as e:
        logger.error(f"Error parsing PPTX {pptx_path}: {e}")
        return ""

def parse_docx_file(docx_path: Path) -> str:
    """
    Parse a DOCX file and extract text content.
    
    Args:
        docx_path: Path to the DOCX file
        
    Returns:
        Extracted text content
    """
    try:
        from docx import Document
        
        text_content = []
        doc = Document(docx_path)
        
        # Extract paragraphs
        for para in doc.paragraphs:
            if para.text.strip():
                text_content.append(para.text.strip())
        
        # Extract tables
        for table_num, table in enumerate(doc.tables, 1):
            text_content.append(f"\n--- Table {table_num} ---")
            for row in table.rows:
                row_text = []
                for cell in row.cells:
                    if cell.text.strip():
                        row_text.append(cell.text.strip())
                if row_text:
                    text_content.append(" | ".join(row_text))
        
        return "\n\n".join(text_content)
        
    except ImportError:
        logger.error("python-docx not installed. Install with: pip install python-docx")
        return ""
    except Exception as e:
        logger.error(f"Error parsing DOCX {docx_path}: {e}")
        return ""

def parse_file_by_type(file_path: Path) -> str:
    """
    Parse a file based on its extension.
    
    Args:
        file_path: Path to the file to parse
        
    Returns:
        Extracted text content
    """
    file_extension = file_path.suffix.lower()
    
    if file_extension == '.pdf':
        return parse_pdf_file(file_path)
    elif file_extension == '.pptx':
        return parse_pptx_file(file_path)
    elif file_extension == '.docx':
        return parse_docx_file(file_path)
    else:
        logger.warning(f"Unsupported file type: {file_extension}")
        return ""

def get_file_metadata(file_path: Path) -> Dict[str, Any]:
    """
    Extract metadata from a file.
    
    Args:
        file_path: Path to the file
        
    Returns:
        Dictionary containing file metadata
    """
    metadata = {
        "file_name": file_path.name,
        "file_size": file_path.stat().st_size,
        "file_extension": file_path.suffix.lower(),
        "file_path": str(file_path)
    }
    
    # Add specific metadata based on file type
    file_extension = file_path.suffix.lower()
    
    if file_extension == '.pdf':
        try:
            import pdfplumber
            with pdfplumber.open(file_path) as pdf:
                metadata.update({
                    "page_count": len(pdf.pages),
                    "title": pdf.metadata.get("Title", "") if pdf.metadata else "",
                    "author": pdf.metadata.get("Author", "") if pdf.metadata else "",
                    "subject": pdf.metadata.get("Subject", "") if pdf.metadata else ""
                })
        except Exception as e:
            logger.warning(f"Could not extract PDF metadata: {e}")
    
    elif file_extension == '.pptx':
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            metadata.update({
                "slide_count": len(prs.slides),
                "title": prs.core_properties.title or "",
                "author": prs.core_properties.author or "",
                "subject": prs.core_properties.subject or ""
            })
        except Exception as e:
            logger.warning(f"Could not extract PPTX metadata: {e}")
    
    elif file_extension == '.docx':
        try:
            from docx import Document
            doc = Document(file_path)
            metadata.update({
                "paragraph_count": len(doc.paragraphs),
                "table_count": len(doc.tables),
                "title": doc.core_properties.title or "",
                "author": doc.core_properties.author or "",
                "subject": doc.core_properties.subject or ""
            })
        except Exception as e:
            logger.warning(f"Could not extract DOCX metadata: {e}")
    
    return metadata
